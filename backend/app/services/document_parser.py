from pathlib import Path

from docx import Document
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_FILE_TYPES = {".txt", ".md", ".pdf", ".docx", ".pptx"}


def normalize_text(text: str) -> str:
    lines = [line.strip() for line in text.replace("\r\n", "\n").split("\n")]
    cleaned_lines: list[str] = []
    previous_blank = False

    for line in lines:
        is_blank = not line
        if is_blank and previous_blank:
            continue
        cleaned_lines.append(line)
        previous_blank = is_blank

    return "\n".join(cleaned_lines).strip()


def parse_document_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix not in SUPPORTED_FILE_TYPES:
        supported = ", ".join(sorted(SUPPORTED_FILE_TYPES))
        raise ValueError(f"Unsupported file type: {suffix}. Supported: {supported}")

    if suffix in {".txt", ".md"}:
        return normalize_text(file_path.read_text(encoding="utf-8"))

    if suffix == ".pdf":
        return normalize_text(_parse_pdf(file_path))

    if suffix == ".docx":
        return normalize_text(_parse_docx(file_path))

    return normalize_text(_parse_pptx(file_path))


def _parse_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    page_texts: list[str] = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            page_texts.append(page_text)
    return "\n\n".join(page_texts)


def _parse_docx(file_path: Path) -> str:
    document = Document(str(file_path))
    parts: list[str] = []

    for paragraph in document.paragraphs:
        text = paragraph.text.strip()
        if text:
            parts.append(text)

    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    return "\n\n".join(parts)


def _parse_pptx(file_path: Path) -> str:
    presentation = Presentation(str(file_path))
    parts: list[str] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                text = shape.text.strip()
                if text:
                    slide_parts.append(text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))

        if slide_parts:
            parts.append(f"Slide {slide_index}\n" + "\n".join(slide_parts))

    return "\n\n".join(parts)
